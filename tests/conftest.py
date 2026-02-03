"""
Pytest fixtures for document extraction tests.

These fixtures generate temporary test documents for use in tests.
"""

import pytest
import fitz  # pymupdf


@pytest.fixture
def tmp_pdf(tmp_path):
    """Generate a simple PDF for testing.

    Creates a PDF with:
    - A title (large font)
    - Two body paragraphs
    - A simple 3x3 table
    - An embedded image (tiny PNG)

    Returns:
        Path to the generated PDF file.
    """
    doc = fitz.open()
    page = doc.new_page()

    # Title (large font - should become H1)
    page.insert_text((72, 72), "Test Document Title", fontsize=18)

    # Subheading (medium font - should become H2)
    page.insert_text((72, 110), "Section One", fontsize=14)

    # Body text (normal font)
    page.insert_text((72, 140), "This is the first body paragraph with some text.", fontsize=11)
    page.insert_text((72, 160), "This is the second body paragraph.", fontsize=11)

    # Create a simple 3x3 table using a Table object
    table_rect = fitz.Rect(72, 200, 300, 290)
    # Draw table grid
    page.draw_rect(table_rect, color=(0, 0, 0), width=0.5)
    # Draw horizontal lines
    for i in range(1, 3):
        y = 200 + i * 30
        page.draw_line((72, y), (300, y), color=(0, 0, 0), width=0.5)
    # Draw vertical lines
    for i in range(1, 3):
        x = 72 + i * 76
        page.draw_line((x, 200), (x, 290), color=(0, 0, 0), width=0.5)

    # Add table cell text
    cells = [
        ["Header1", "Header2", "Header3"],
        ["Row1Col1", "Row1Col2", "Row1Col3"],
        ["Row2Col1", "Row2Col2", "Row2Col3"],
    ]
    for row_idx, row in enumerate(cells):
        for col_idx, cell in enumerate(row):
            x = 75 + col_idx * 76
            y = 220 + row_idx * 30
            page.insert_text((x, y), cell, fontsize=9)

    # Create a tiny PNG image (1x1 red pixel)
    import io
    import struct
    import zlib

    def create_minimal_png():
        """Create a minimal valid PNG (1x1 red pixel)."""
        # PNG signature
        signature = b'\x89PNG\r\n\x1a\n'
        # IHDR chunk
        ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
        ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
        ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
        # IDAT chunk (red pixel)
        raw_data = b'\x00\xff\x00\x00'  # filter byte + RGB
        compressed = zlib.compress(raw_data)
        idat_crc = zlib.crc32(b'IDAT' + compressed) & 0xffffffff
        idat = struct.pack('>I', len(compressed)) + b'IDAT' + compressed + struct.pack('>I', idat_crc)
        # IEND chunk
        iend_crc = zlib.crc32(b'IEND') & 0xffffffff
        iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
        return signature + ihdr + idat + iend

    png_data = create_minimal_png()
    img_rect = fitz.Rect(72, 310, 122, 360)
    page.insert_image(img_rect, stream=png_data)

    path = tmp_path / "test.pdf"
    doc.save(path)
    doc.close()
    return path


@pytest.fixture
def tmp_docx(tmp_path):
    """Generate a simple DOCX for testing.

    Creates a DOCX with:
    - A Heading 1
    - A Heading 2
    - Body paragraphs
    - A 2x3 table

    Returns:
        Path to the generated DOCX file.
    """
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document Title", level=1)
    doc.add_heading("Section One", level=2)
    doc.add_paragraph("This is the first body paragraph.")
    doc.add_paragraph("This is the second body paragraph.")

    # Add a 2x3 table
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Header1"
    table.cell(0, 1).text = "Header2"
    table.cell(0, 2).text = "Header3"
    table.cell(1, 0).text = "Data1"
    table.cell(1, 1).text = "Data2"
    table.cell(1, 2).text = "Data3"

    path = tmp_path / "test.docx"
    doc.save(path)
    return path


@pytest.fixture
def tmp_pptx(tmp_path):
    """Generate a minimal PPTX for testing.

    Creates a PPTX with one slide containing a title.
    Needed to instantiate PPTXExtractor even though extraction
    isn't implemented.

    Returns:
        Path to the generated PPTX file.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Test Presentation"

    path = tmp_path / "test.pptx"
    prs.save(path)
    return path


@pytest.fixture
def tmp_xlsx(tmp_path):
    """Generate a minimal XLSX for testing.

    Creates an XLSX with one sheet containing a few cells.
    Needed to instantiate XLSXExtractor even though extraction
    isn't implemented.

    Returns:
        Path to the generated XLSX file.
    """
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws['A1'] = 'Header1'
    ws['B1'] = 'Header2'
    ws['A2'] = 'Data1'
    ws['B2'] = 'Data2'

    path = tmp_path / "test.xlsx"
    wb.save(path)
    return path


@pytest.fixture
def tmp_empty_pdf(tmp_path):
    """Generate an empty PDF for edge case testing.

    Creates a PDF with no content (just a blank page).

    Returns:
        Path to the generated PDF file.
    """
    doc = fitz.open()
    doc.new_page()  # Blank page
    path = tmp_path / "empty.pdf"
    doc.save(path)
    doc.close()
    return path


@pytest.fixture
def sample_dir(tmp_path):
    """Return a temporary directory for test outputs.

    Returns:
        Path to temporary directory.
    """
    return tmp_path
